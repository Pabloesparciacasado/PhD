# In[]:
import pandas as pd
import numpy as np

import io
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


opt_df = pd.read_parquet(r"C:\Users\pablo.esparcia\Documents\OptionMetrics\output\preliminares\opt_df_prueba.parquet")

generar_pptx = True   # ← cambiar a True para exportar el PowerPoint

_slides_buf = []   # buffers PNG para el PowerPoint

def _save_slide(title_text=""):
    """Guarda la figura en el buffer y la cierra."""
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    _slides_buf.append((title_text, buf))
    plt.close()

def _fmt(x, _):
    """Formatter adaptativo: B / M / K según magnitud, siempre positivo."""
    v = abs(x)
    if v >= 1e9:   return f"{v/1e9:.1f}B"
    if v >= 1e6:   return f"{v/1e6:.1f}M"
    if v >= 1e3:   return f"{v/1e3:.1f}K"
    if v > 0:      return f"{v:.2f}"
    return "0"



##### visualicemos la distribución del volumen por moneyness:
# In[]:

df_plot = opt_df[
    opt_df["Volume"].notna() &
    # (opt_df["Volume"] >= 0) &
    opt_df["Moneyness"].notna()
].copy()

bins = np.linspace(0.6, 1.4, 100)  # moneyness de 0.6 a 1.4

calls = df_plot[df_plot["CallPut"] == "C"]
puts  = df_plot[df_plot["CallPut"] == "P"]

vol_calls, _ = np.histogram(calls["Moneyness"], bins=bins, weights=calls["Volume"])
vol_puts,  _ = np.histogram(puts["Moneyness"],  bins=bins, weights=puts["Volume"])

bin_centers = (bins[:-1] + bins[1:]) / 2
width = bins[1] - bins[0]


fig, ax = plt.subplots(figsize=(12, 5))

ax.bar(bin_centers, vol_calls / 1e6, width=width * 0.9,
       label="Calls", color="#2196F3", alpha=0.8)
ax.bar(bin_centers, -vol_puts / 1e6, width=width * 0.9,
       label="Puts",  color="#F44336", alpha=0.8)

ax.axvline(1, color="black", linewidth=1, linestyle="--", label="ATM (M = 1)")
ax.axhline(0, color="black", linewidth=0.8)

ax.set_xlabel("Moneyness  [Strike / Spot]", fontsize=12)
ax.set_ylabel("Volumen agregado (millones de contratos)", fontsize=12)
ax.set_title("Distribución de volumen por moneyness — Calls  vs Puts ", fontsize=13)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(_fmt))
ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=20, symmetric=True))
ax.legend()
plt.tight_layout()
_save_slide("Volumen — Total")

# In[]:
##### distribución del volumen por moneyness — desglose por bucket de vencimiento:

expiry_buckets = [
    (0,   15,  "< 15 días"),
    (15,  45,  "15–45 días"),
    (45,  105, "45–105 días"),
    (105, 183, "105–183 días"),
    (183, 365, "183–365 días"),
    (365, np.inf, "≥ 365 días"),
]

mono_bins = np.linspace(0.6, 1.4, 100)
bin_centers_m = (mono_bins[:-1] + mono_bins[1:]) / 2
width_m = mono_bins[1] - mono_bins[0]

for lo, hi, label in expiry_buckets:
    mask = (df_plot["Days"] >= lo) & (df_plot["Days"] < hi)
    sub = df_plot[mask]

    calls_s = sub[sub["CallPut"] == "C"]
    puts_s  = sub[sub["CallPut"] == "P"]

    vc, _ = np.histogram(calls_s["Moneyness"], bins=mono_bins, weights=calls_s["Volume"])
    vp, _ = np.histogram(puts_s["Moneyness"],  bins=mono_bins, weights=puts_s["Volume"])

    fig, ax = plt.subplots(figsize=(12, 5))
    ax.bar(bin_centers_m,  vc / 1e6, width=width_m * 0.9, label="Calls", color="#2196F3", alpha=0.8)
    ax.bar(bin_centers_m, -vp / 1e6, width=width_m * 0.9, label="Puts",  color="#F44336", alpha=0.8)
    ax.axvline(1, color="black", linewidth=1, linestyle="--", label="ATM (M = 1)")
    ax.axhline(0, color="black", linewidth=0.6)
    ax.set_title(f"Distribución de volumen por moneyness — {label}", fontsize=13)
    ax.set_xlabel("Moneyness  [Strike / Spot]", fontsize=12)
    ax.set_ylabel("Volumen agregado (millones de contratos)", fontsize=12)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(_fmt))
    ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=20, symmetric=True))
    ax.legend()
    plt.tight_layout()
    _save_slide(f"Volumen — {label}")

# In[]:
##### distribución del Open Interest por moneyness:

df_plot_oi = opt_df[
    opt_df["OpenInterest"].notna() &
    # (opt_df["OpenInterest"] > 0) &
    opt_df["Moneyness"].notna()
].copy()

calls_oi = df_plot_oi[df_plot_oi["CallPut"] == "C"]
puts_oi  = df_plot_oi[df_plot_oi["CallPut"] == "P"]

oi_calls, _ = np.histogram(calls_oi["Moneyness"], bins=bins, weights=calls_oi["OpenInterest"])
oi_puts,  _ = np.histogram(puts_oi["Moneyness"],  bins=bins, weights=puts_oi["OpenInterest"])

fig, ax = plt.subplots(figsize=(12, 5))
ax.bar(bin_centers,  oi_calls / 1e6, width=width * 0.9, label="Calls", color="#2196F3", alpha=0.8)
ax.bar(bin_centers, -oi_puts  / 1e6, width=width * 0.9, label="Puts",  color="#F44336", alpha=0.8)
ax.axvline(1, color="black", linewidth=1, linestyle="--", label="ATM (M = 1)")
ax.axhline(0, color="black", linewidth=0.8)
ax.set_xlabel("Moneyness  [Strike / Spot]", fontsize=12)
ax.set_ylabel("Open Interest agregado (millones de contratos)", fontsize=12)
ax.set_title("Distribución de Open Interest por moneyness — Calls  vs Puts ", fontsize=13)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(_fmt))
ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=20, symmetric=True))
ax.legend()
plt.tight_layout()
_save_slide("Open Interest — Total")

# In[]:
##### distribución del Open Interest por moneyness — desglose por bucket de vencimiento:

for lo, hi, label in expiry_buckets:
    mask = (df_plot_oi["Days"] >= lo) & (df_plot_oi["Days"] < hi)
    sub = df_plot_oi[mask]

    calls_s = sub[sub["CallPut"] == "C"]
    puts_s  = sub[sub["CallPut"] == "P"]

    vc, _ = np.histogram(calls_s["Moneyness"], bins=mono_bins, weights=calls_s["OpenInterest"])
    vp, _ = np.histogram(puts_s["Moneyness"],  bins=mono_bins, weights=puts_s["OpenInterest"])

    fig, ax = plt.subplots(figsize=(12, 5))
    ax.bar(bin_centers_m,  vc / 1e6, width=width_m * 0.9, label="Calls", color="#2196F3", alpha=0.8)
    ax.bar(bin_centers_m, -vp / 1e6, width=width_m * 0.9, label="Puts",  color="#F44336", alpha=0.8)
    ax.axvline(1, color="black", linewidth=1, linestyle="--", label="ATM (M = 1)")
    ax.axhline(0, color="black", linewidth=0.6)
    ax.set_title(f"Distribución de Open Interest por moneyness — {label}", fontsize=13)
    ax.set_xlabel("Moneyness  [Strike / Spot]", fontsize=12)
    ax.set_ylabel("Open Interest agregado (millones de contratos)", fontsize=12)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(_fmt))
    ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=20, symmetric=True))
    ax.legend()
    plt.tight_layout()
    _save_slide(f"Open Interest — {label}")

# In[]:
##### distribución del Dollar Volume por moneyness:

opt_df["DollarVolume"] = opt_df["Volume"] * opt_df["MidPrice"]

df_plot_dv = opt_df[
    opt_df["DollarVolume"].notna() &
    # (opt_df["DollarVolume"] > 0) &
    opt_df["Moneyness"].notna()
].copy()

calls_dv = df_plot_dv[df_plot_dv["CallPut"] == "C"]
puts_dv  = df_plot_dv[df_plot_dv["CallPut"] == "P"]

dv_calls, _ = np.histogram(calls_dv["Moneyness"], bins=bins, weights=calls_dv["DollarVolume"])
dv_puts,  _ = np.histogram(puts_dv["Moneyness"],  bins=bins, weights=puts_dv["DollarVolume"])

fig, ax = plt.subplots(figsize=(12, 5))
ax.bar(bin_centers,  dv_calls / 1e9, width=width * 0.9, label="Calls", color="#2196F3", alpha=0.8)
ax.bar(bin_centers, -dv_puts  / 1e9, width=width * 0.9, label="Puts",  color="#F44336", alpha=0.8)
ax.axvline(1, color="black", linewidth=1, linestyle="--", label="ATM (M = 1)")
ax.axhline(0, color="black", linewidth=0.8)
ax.set_xlabel("Moneyness  [Strike / Spot]", fontsize=12)
ax.set_ylabel("Dollar Volume agregado (miles de millones $)", fontsize=12)
ax.set_title("Distribución de Dollar Volume por moneyness — Calls  vs Puts ", fontsize=13)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(_fmt))
ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=20, symmetric=True))
ax.legend()
plt.tight_layout()
_save_slide("Dollar Volume — Total")

# In[]:
##### distribución del Dollar Volume por moneyness — desglose por bucket de vencimiento:

for lo, hi, label in expiry_buckets:
    mask = (df_plot_dv["Days"] >= lo) & (df_plot_dv["Days"] < hi)
    sub = df_plot_dv[mask]

    calls_s = sub[sub["CallPut"] == "C"]
    puts_s  = sub[sub["CallPut"] == "P"]

    vc, _ = np.histogram(calls_s["Moneyness"], bins=mono_bins, weights=calls_s["DollarVolume"])
    vp, _ = np.histogram(puts_s["Moneyness"],  bins=mono_bins, weights=puts_s["DollarVolume"])

    fig, ax = plt.subplots(figsize=(12, 5))
    ax.bar(bin_centers_m,  vc / 1e9, width=width_m * 0.9, label="Calls", color="#2196F3", alpha=0.8)
    ax.bar(bin_centers_m, -vp / 1e9, width=width_m * 0.9, label="Puts",  color="#F44336", alpha=0.8)
    ax.axvline(1, color="black", linewidth=1, linestyle="--", label="ATM (M = 1)")
    ax.axhline(0, color="black", linewidth=0.6)
    ax.set_title(f"Distribución de Dollar Volume por moneyness — {label}", fontsize=13)
    ax.set_xlabel("Moneyness  [Strike / Spot]", fontsize=12)
    ax.set_ylabel("Dollar Volume agregado (miles de millones $)", fontsize=12)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(_fmt))
    ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=20, symmetric=True))
    ax.legend()
    plt.tight_layout()
    _save_slide(f"Dollar Volume — {label}")

# In[]:
##### cobertura de datos diaria por bucket de vencimiento y moneyness:

# Columna de fecha — ajustar si el nombre es diferente
DATE_COL = "Date"  
# Buckets de moneyness más gruesos para la matriz de cobertura
mono_cov_bins   = np.array([0.3, 0.4, 0.5,0.6, 0.7, 0.8, 0.9, 1.0, 1.1, 1.2, 1.3, 1.4])
mono_cov_labels = ["0.3–0.4", "0.4–0.5", "0.5–0.6", "0.6–0.7", "0.7–0.8", "0.8–0.9", "0.9–1.0",
                   "1.0–1.1", "1.1–1.2", "1.2–1.3", "1.3–1.4"]
exp_labels = [lbl for _, _, lbl in expiry_buckets]

def _assign_expiry_bucket(days):
    for lo, hi, lbl in expiry_buckets:
        if lo <= days < hi:
            return lbl
    return None

def _build_coverage(df, metric_col):
    """
    Devuelve (cov_pct, days_mat, total_days):
      - cov_pct   : DataFrame (expiry x mono) con % de dias que tienen metric_col > 0
      - days_mat  : idem con numero absoluto de dias
      - total_days: dias unicos en la muestra
    """
    df = df.copy()
    df["_mono"] = pd.cut(df["Moneyness"], bins=mono_cov_bins,
                         labels=mono_cov_labels, right=False)
    df["_exp"]  = df["Days"].apply(_assign_expiry_bucket)
    df = df.dropna(subset=["_mono", "_exp", metric_col])
    df = df[df[metric_col] > 0]

    total_days = df[DATE_COL].nunique()

    # para cada combinacion (date, expiry, mono) basta con saber si hay al menos 1 fila
    presence = (
        df.groupby([DATE_COL, "_exp", "_mono"], observed=True)
          .size()
          .reset_index(name="_n")
    )
    days_mat_raw = (
        presence.groupby(["_exp", "_mono"], observed=True)
                .size()
                .unstack("_mono", fill_value=0)
    )
    days_mat = days_mat_raw.reindex(index=exp_labels, columns=mono_cov_labels, fill_value=0)
    cov_pct  = days_mat / total_days * 100 if total_days > 0 else days_mat * 0
    return cov_pct, days_mat, total_days

def _plot_coverage(cov_pct, days_mat, total_days, title, slide_title):
    fig, ax = plt.subplots(figsize=(12, 4))
    im = ax.imshow(cov_pct.values, aspect="auto", cmap="YlGn", vmin=0, vmax=100)

    ax.set_xticks(range(len(mono_cov_labels)))
    ax.set_xticklabels(mono_cov_labels, rotation=45, ha="right", fontsize=10)
    ax.set_yticks(range(len(exp_labels)))
    ax.set_yticklabels(exp_labels, fontsize=10)

    for i in range(len(exp_labels)):
        for j in range(len(mono_cov_labels)):
            pct = cov_pct.values[i, j]
            n   = int(days_mat.values[i, j])
            txt_color = "white" if pct >= 65 else "black"
            ax.text(j, i, f"{pct:.0f}%\n({n}d)",
                    ha="center", va="center", fontsize=8, color=txt_color)

    cbar = plt.colorbar(im, ax=ax, fraction=0.03, pad=0.04)
    cbar.set_label("% dias con datos", fontsize=10)
    ax.set_xlabel("Moneyness bucket", fontsize=11)
    ax.set_ylabel("Bucket de vencimiento", fontsize=11)
    ax.set_title(f"{title}  (total muestra: {total_days} dias)", fontsize=13)
    plt.tight_layout()
    _save_slide(slide_title)

# Volumen — Calls / Puts
for cp, cp_lbl in [("C", "Calls"), ("P", "Puts")]:
    sub = opt_df[
        opt_df["Volume"].notna() &
        opt_df["Moneyness"].notna() &
        (opt_df["CallPut"] == cp)
    ]
    cov_pct, days_mat, total_days = _build_coverage(sub, "Volume")
    _plot_coverage(cov_pct, days_mat, total_days,
                   title=f"Cobertura de Volumen — {cp_lbl}",
                   slide_title=f"Cobertura Volumen — {cp_lbl}")

# Open Interest — Calls / Puts
for cp, cp_lbl in [("C", "Calls"), ("P", "Puts")]:
    sub = opt_df[
        opt_df["OpenInterest"].notna() &
        opt_df["Moneyness"].notna() &
        (opt_df["CallPut"] == cp)
    ]
    cov_pct, days_mat, total_days = _build_coverage(sub, "OpenInterest")
    _plot_coverage(cov_pct, days_mat, total_days,
                   title=f"Cobertura de Open Interest — {cp_lbl}",
                   slide_title=f"Cobertura OI — {cp_lbl}")

# In[]:
##### número de contratos activos (Vol>0 o OI>0) por bucket de vencimiento y moneyness:

df_active = opt_df[
    opt_df["Moneyness"].notna() &
    ((opt_df["Volume"] > 0) | (opt_df["OpenInterest"] > 0))
].copy()
df_active["_mono"] = pd.cut(df_active["Moneyness"], bins=mono_cov_bins,
                             labels=mono_cov_labels, right=False)
df_active["_exp"]  = df_active["Days"].apply(_assign_expiry_bucket)
df_active = df_active.dropna(subset=["_mono", "_exp"])

def _plot_contracts(df, cp, cp_lbl):
    sub = df[df["CallPut"] == cp]
    counts_raw = (
        sub.groupby(["_exp", "_mono"], observed=True)
           .size()
           .unstack("_mono", fill_value=0)
    )
    counts = counts_raw.reindex(index=exp_labels, columns=mono_cov_labels, fill_value=0)

    fig, ax = plt.subplots(figsize=(12, 4))
    # escala logarítmica para que los extremos no queden blancos
    import matplotlib.colors as mcolors
    vals = counts.values.astype(float)
    vmin = max(vals[vals > 0].min(), 1) if (vals > 0).any() else 1
    norm = mcolors.LogNorm(vmin=vmin, vmax=vals.max() + 1)
    im = ax.imshow(vals, aspect="auto", cmap="YlGn", norm=norm)

    ax.set_xticks(range(len(mono_cov_labels)))
    ax.set_xticklabels(mono_cov_labels, rotation=45, ha="right", fontsize=10)
    ax.set_yticks(range(len(exp_labels)))
    ax.set_yticklabels(exp_labels, fontsize=10)

    for i in range(len(exp_labels)):
        for j in range(len(mono_cov_labels)):
            n = int(vals[i, j])
            # umbral de color basado en posición normalizada en log-scale
            normed = (np.log1p(n) - np.log1p(vmin)) / (np.log1p(vals.max()) - np.log1p(vmin) + 1e-9)
            txt_color = "white" if normed >= 0.65 else "black"
            label_n = _fmt(n, None) if n > 0 else "0"
            ax.text(j, i, label_n, ha="center", va="center",
                    fontsize=8, color=txt_color)

    cbar = plt.colorbar(im, ax=ax, fraction=0.03, pad=0.04)
    cbar.set_label("Nº contratos (escala log)", fontsize=10)
    ax.set_xlabel("Moneyness bucket", fontsize=11)
    ax.set_ylabel("Bucket de vencimiento", fontsize=11)
    total = int(sub.shape[0])
    ax.set_title(
        f"Contratos con Vol>0 o OI>0 — {cp_lbl}  (total: {_fmt(total, None)} contratos)",
        fontsize=13
    )
    plt.tight_layout()
    _save_slide(f"Nº Contratos — {cp_lbl}")

for cp, cp_lbl in [("C", "Calls"), ("P", "Puts")]:
    _plot_contracts(df_active, cp, cp_lbl)

# In[]:
##### % de días con al menos un contrato activo (Vol>0 o OI>0) por bucket:

for cp, cp_lbl in [("C", "Calls"), ("P", "Puts")]:
    sub = df_active[df_active["CallPut"] == cp]
    total_days = sub[DATE_COL].nunique()

    presence = (
        sub.groupby([DATE_COL, "_exp", "_mono"], observed=True)
           .size()
           .reset_index(name="_n")
    )
    days_mat_raw = (
        presence.groupby(["_exp", "_mono"], observed=True)
                .size()
                .unstack("_mono", fill_value=0)
    )
    days_mat = days_mat_raw.reindex(index=exp_labels, columns=mono_cov_labels, fill_value=0)
    cov_pct  = days_mat / total_days * 100 if total_days > 0 else days_mat * 0

    _plot_coverage(cov_pct, days_mat, total_days,
                   title=f"Cobertura de días (Vol>0 o OI>0) — {cp_lbl}",
                   slide_title=f"Cobertura días activos — {cp_lbl}")

# In[]:
##### exportar todos los gráficos a PowerPoint:
if generar_pptx:

    PPTX_OUT = r"C:\Users\pablo.esparcia\Documents\OptionMetrics\output\preliminares\moneyness_distribution.pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]   # layout en blanco

    # --- Diapositiva de índice ---
    idx_slide = prs.slides.add_slide(blank_layout)

    # cabecera azul (igual que el resto)
    hdr = idx_slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.55)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
    hdr.line.fill.background()

    txb_hdr = idx_slide.shapes.add_textbox(Inches(0.2), Inches(0.05), Inches(12.9), Inches(0.45))
    tf_hdr = txb_hdr.text_frame
    tf_hdr.text = "Índice"
    tf_hdr.paragraphs[0].runs[0].font.size  = Pt(18)
    tf_hdr.paragraphs[0].runs[0].font.bold  = True
    tf_hdr.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # lista de slides en dos columnas
    titles = [t for t, _ in _slides_buf]
    mid    = (len(titles) + 1) // 2
    col_left  = titles[:mid]
    col_right = titles[mid:]

    for col_idx, col_items in enumerate([col_left, col_right]):
        x_pos = Inches(0.4 + col_idx * 6.4)
        txb = idx_slide.shapes.add_textbox(x_pos, Inches(0.7), Inches(6.2), Inches(6.6))
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, title in enumerate(col_items, start=1 + col_idx * mid):
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = f"{i}.  {title}"
            p.runs[0].font.size  = Pt(13)
            p.runs[0].font.color.rgb = RGBColor(0x1E, 0x27, 0x61)
            p.space_after = Pt(4)

    for slide_title, buf in _slides_buf:
        slide = prs.slides.add_slide(blank_layout)

        header_shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), prs.slide_width, Inches(0.55)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
        header_shape.line.fill.background()

        txb = slide.shapes.add_textbox(Inches(0.2), Inches(0.05), Inches(12.9), Inches(0.45))
        tf = txb.text_frame
        tf.text = slide_title
        tf.paragraphs[0].runs[0].font.size  = Pt(18)
        tf.paragraphs[0].runs[0].font.bold  = True
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        slide.shapes.add_picture(buf, Inches(0), Inches(0.55),
                                 width=prs.slide_width,
                                 height=prs.slide_height - Inches(0.55))

    prs.save(PPTX_OUT)
    print(f"PowerPoint guardado en: {PPTX_OUT}")

#####
# Análisis díario:
